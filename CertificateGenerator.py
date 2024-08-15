import pandas as pd
from pptx import Presentation
from pptx.util import Pt
import os

# Load the Excel file
df = pd.read_excel('C://Users//EXCELWITHNAMES.xlsx')

# Ensure the columns are correctly read
df['Participantes'] = df['Participantes'].astype(str)
df['Cédula'] = df['Cédula '].astype(str)
df['Tema'] = df['Tema'].astype(str)

# Load the PowerPoint template
template_path = 'C://Users//POWERPOINTEMPLATE.pptx'
output_folder = 'C://Users//GeneratedFiles'

# Ensure the output folder exists
if not os.path.exists(output_folder):
    os.makedirs(output_folder)

# Group by 'Participantes' to collect all topics per participant
grouped = df.groupby(['Participantes', 'Cédula'])['Tema'].apply(lambda x: '\n'.join(x)).reset_index()

# Loop over each unique participant and create a PowerPoint file
for index, row in grouped.iterrows():
    participant_name = row['Participantes']
    participant_id = row['Cédula']
    temas = '-' + row['Tema']
    
    # Open the PowerPoint template
    prs = Presentation(template_path)
    
    # Edit text placeholders in the presentation
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                
                # Replace placeholder for the participant's name
                if 'JUST PASTE WHAT THE FIRST TEXT BOX YOU WANT TO CHANGE SAYS' in shape.text:
                    shape.text_frame.clear()
                    p = shape.text_frame.add_paragraph()
                    p.text = participant_name
                    p.font.name = 'Gill Sans MT'
                    p.font.size = Pt(24)
                    p.font.bold = True
                    
                # Replace placeholder for the participant's ID
                elif 'JUST PASTE WHAT THE SECOND TEXT BOX YOU WANT TO CHANGE SAYS' in shape.text:
                    shape.text_frame.clear()
                    p = shape.text_frame.add_paragraph()
                    p.text = f"C.C. {participant_id}"
                    p.font.name = 'Gill Sans MT'
                    p.font.size = Pt(24)
                    p.font.bold = True
                
                # Replace placeholder for the topics (Tema)
                elif 'JUST PASTE WHAT THE THIRD TEXT BOX YOU WANT TO CHANGE SAYS' in shape.text:
                    shape.text_frame.clear()
                    p = shape.text_frame.add_paragraph()
                    p.text = temas
                    p.font.name = 'Gill Sans MT'
                    p.font.size = Pt(14)
                    p.font.bold = False
    
    # Save the modified PowerPoint for this participant
    pptx_filename = os.path.join(output_folder, f"{participant_name}.pptx")
    prs.save(pptx_filename)

print("PowerPoint files have been generated successfully.")
